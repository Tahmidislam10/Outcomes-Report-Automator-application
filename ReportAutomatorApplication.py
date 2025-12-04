from tkinter import *
from tkinter import messagebox
from PIL import Image, ImageTk
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from io import BytesIO

# Global variables
regions = ['Central London']  # Default region selection
intervention_dates = ['created', 'date_of_agreement', 'end_date']

# Main window setup
window = Tk()
window.geometry("940x620")
window.title("Outcomes Report Automator Application")


# Function to clear the window
def clear_window():
    for widget in window.winfo_children():
        widget.destroy()


# Function to create the background canvas with an image
def create_background():
    global canvas, bg_image
    canvas = Canvas(window, width=940, height=620)
    canvas.pack(fill="both", expand=True)

    # Load and resize the background image
    bg_image = Image.open("stars.png")
    bg_image = bg_image.resize((940, 620), Image.LANCZOS)
    bg_image = ImageTk.PhotoImage(bg_image)

    # Display the background image
    canvas.create_image(0, 0, anchor=NW, image=bg_image)


# Function to add a corner image
def add_corner_image():
    corner_image = Image.open("Place2be.png")
    corner_image = corner_image.resize((100, 75), Image.LANCZOS)
    corner_image = ImageTk.PhotoImage(corner_image)

    corner_label = Label(window, image=corner_image, background='orange')
    corner_label.image = corner_image
    corner_label.place(x=820, y=10)


# Function to add an info label
def add_info_label():
    info_label = Label(window, text="ⓘResearch and Evaluation Team", font=('Arial', 14), fg='blue', background='orange')
    info_label.place(x=10, y=580)


# Function to handle login and display report selection
def start():
    username = entry1.get()
    password = entry2.get()

    if username == '' and password == '':
        messagebox.showerror('Login', 'Blanks are not allowed!')
    elif username in ['Tahmid Islam', 'Amy Capel', 'Nikhil Naag', 'Hizkia Yolanda'] and password == '12345':
        messagebox.showinfo('Login', 'Login is Successful')
        show_report_selection()
    else:
        messagebox.showerror('Login', 'Invalid username or password')


# Function to exit the application
def exit_app():
    window.destroy()


# Function to show report selection options
def show_report_selection():
    clear_window()
    create_background()

    label2 = Label(window, text="Which Report would you like to select?", font=('Amasis MT Pro Black', 30, 'bold'),
                   background='orange', bd=10)
    label2.place(x=40, y=200)

    button2 = Button(window, text="Area Report", command=area_report_message, font=('Amasis MT Pro Black', 20, 'bold'),
                     fg='green', background='orange', relief=SUNKEN, bd=10)
    button2.place(x=220, y=285)

    button2 = Button(window, text="Region Report", command=what_region, font=('Amasis MT Pro Black', 20, 'bold'),
                     fg='green', background='orange', relief=SUNKEN, bd=10)
    button2.place(x=500, y=285)

    button = Button(window, text="EXIT", command=exit_app, font=('Amasis MT Pro Black', 20, 'bold'), fg='red',
                    background='orange', relief=SUNKEN, bd=10)
    button.place(x=755, y=500)

    add_corner_image()
    add_info_label()


# Function for area report message
def area_report_message():
    messagebox.showinfo("Area Report", "This section is still under development by the research and evaluation team.")


# Function to handle login page display
def login():
    clear_window()
    create_background()

    label2 = Label(window, text="Please Enter Details to login!", font=('Amasis MT Pro Black', 30, 'bold'),
                   background='orange', bd=10)
    label2.place(x=40, y=70)

    label3 = Label(window, text="Username:", font=('Amasis MT Pro Black', 20, 'bold'),
                   background='orange', bd=10)
    label3.place(x=75, y=195)
    global entry1
    entry1 = Entry(window, font=("Arial", 20))
    entry1.place(x=275, y=200)

    label4 = Label(window, text="Password:", font=('Amasis MT Pro Black', 20, 'bold'),
                   background='orange', bd=10)
    label4.place(x=75, y=300)
    global entry2
    entry2 = Entry(window, font=("Arial", 20), show='*')
    entry2.place(x=275, y=300)

    enter_button = Button(window, text="Sign in", bg='cyan', font=('Arial', 20), command=start)
    enter_button.place(x=675, y=200)

    button = Button(window, text="EXIT", command=exit_app, font=('Amasis MT Pro Black', 20, 'bold'), fg='red',
                    background='orange', relief=SUNKEN, bd=10)
    button.place(x=755, y=500)

    add_corner_image()
    add_info_label()


# Function for region selection
def what_region():
    clear_window()
    create_background()

    label2 = Label(window, text="What Area would you like to Select?", font=('Amasis MT Pro Black', 30, 'bold'),
                   background='orange', bd=10)
    label2.place(x=40, y=200)

    var1 = IntVar()
    var2 = IntVar()

    checkbutton1 = Checkbutton(window, text="Central London", variable=var1, font=('Amasis MT Pro Black', 20, 'bold'),
                               fg='green', background='orange', relief=SUNKEN, bd=10)
    checkbutton1.place(x=80, y=285)

    checkbutton2 = Checkbutton(window, text="North West", variable=var2, font=('Amasis MT Pro Black', 20, 'bold'),
                               fg='green', background='orange', relief=SUNKEN, bd=10)
    checkbutton2.place(x=420, y=285)

    button = Button(window, text="EXIT", command=exit_app, font=('Amasis MT Pro Black', 20, 'bold'), fg='red',
                    background='orange', relief=SUNKEN, bd=10)
    button.place(x=755, y=500)

    next_button = Button(window, text="Next", command=region, font=('Amasis MT Pro Black', 20, 'bold'), fg='blue',
                         background='orange', relief=SUNKEN, bd=10)
    next_button.place(x=500, y=400)

    add_corner_image()
    add_info_label()


# Function to handle region selection
def region():
    clear_window()
    create_background()

    label2 = Label(window, text="Which type of report would you like to select?",
                   font=('Amasis MT Pro Black', 30, 'bold'), background='orange', bd=10)
    label2.place(x=40, y=200)

    var1 = IntVar()
    var2 = IntVar()
    var3 = IntVar()

    checkbutton1 = Checkbutton(window, text="SDQ Baseline", variable=var1, font=('Amasis MT Pro Black', 20, 'bold'),
                               fg='green', background='orange', relief=SUNKEN, bd=10)
    checkbutton1.place(x=80, y=285)

    checkbutton2 = Checkbutton(window, text="SDQ Improvements", variable=var2, font=('Amasis MT Pro Black', 20, 'bold'),
                               fg='green', background='orange', relief=SUNKEN, bd=10)
    checkbutton2.place(x=320, y=285)

    checkbutton3 = Checkbutton(window, text="Demographics", variable=var3, font=('Amasis MT Pro Black', 20, 'bold'),
                               fg='green', background='orange', relief=SUNKEN, bd=10)
    checkbutton3.place(x=630, y=285)

    next_button = Button(window, text="Next", command=Demographics, font=('Amasis MT Pro Black', 20, 'bold'), fg='blue',
                         background='orange', relief=SUNKEN, bd=10)
    next_button.place(x=500, y=400)

    button = Button(window, text="EXIT", command=exit_app, font=('Amasis MT Pro Black', 20, 'bold'), fg='red',
                    background='orange', relief=SUNKEN, bd=10)
    button.place(x=755, y=500)

    add_corner_image()
    add_info_label()


# Function for demographics report selection
def Demographics():
    clear_window()
    create_background()

    label2 = Label(window, text="Which type of report would you like to select?",
                   font=('Amasis MT Pro Black', 30, 'bold'), background='orange', bd=10)
    label2.place(x=40, y=200)

    button_excel = Button(window, text="Output Report as Excel", command=create_excel_report,
                          font=('Amasis MT Pro Black', 20, 'bold'), fg='green',
                          background='orange', relief=SUNKEN, bd=10)
    button_excel.place(x=80, y=285)

    button_ppt = Button(window, text="Output Report as PowerPoint", command=create_powerpoint_report,
                        font=('Amasis MT Pro Black', 20, 'bold'), fg='green',
                        background='orange', relief=SUNKEN, bd=10)
    button_ppt.place(x=450, y=285)

    button = Button(window, text="EXIT", command=exit_app, font=('Amasis MT Pro Black', 20, 'bold'), fg='red',
                    background='orange', relief=SUNKEN, bd=10)
    button.place(x=755, y=500)

    add_corner_image()
    add_info_label()


# Function to create Excel report
def create_excel_report():
    try:
        # Read data from CSV files (replace with actual file paths)
        t_intervention = pd.read_csv('t_intervention.csv', parse_dates=intervention_dates, dayfirst=True)
        t_quest_SDQ = pd.read_csv('t_quest_SDQ.csv')
        t_lookup_item = pd.read_csv('t_lookup_item.csv')
        t_service_area = pd.read_csv('t_service_area.csv')
        t_cluster = pd.read_csv('t_cluster.csv')
        t_location = pd.read_csv('t_location.csv')
        t_intervention_timeline_task = pd.read_csv('t_intervention_timeline_task.csv')
        t_timeline_task = pd.read_csv('t_timeline_task.csv')

        # Data processing (merge, filter, etc.) - Replace with actual data processing logic
        # For example, merging dataframes
        merged_data = pd.merge(t_intervention, t_quest_SDQ, on='common_column')

        # Example: Filter data based on conditions
        filtered_data = merged_data[merged_data['some_column'] > 100]

        # Save filtered data to Excel
        filtered_data.to_excel('demographics_report.xlsx', index=False)

        # Show success message
        messagebox.showinfo("Excel Report", "Excel report generated successfully!")

    except FileNotFoundError:
        messagebox.showerror("File Error", "One or more CSV files not found.")

    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {str(e)}")


from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO

# Path to the provided template
template_path = 'C:\Outcomes Report Automator application\Outcomes Report automator\Template.pptx'


# Function to create a new PowerPoint report using the given template
def create_powerpoint_report():
    try:
        # Initialize a presentation object
        prs = Presentation()

        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "London and South Data Request"
        subtitle.text = "Research and Evaluation Team"

        # Slide 2: Summary
        summary_slide_layout = prs.slide_layouts[1]
        summary_slide = prs.slides.add_slide(summary_slide_layout)
        title = summary_slide.shapes.title
        content = summary_slide.placeholders[1]
        title.text = "Summary"
        content.text = (
            "This presentation provides a comprehensive analysis of data "
            "from the London and South regions, focusing on key metrics and outcomes."
        )

        # Add more slides as per the structure of your provided PowerPoint
        # For demonstration, I'll add a couple more slides with different layouts

        # Slide 3: Data Overview
        data_overview_layout = prs.slide_layouts[1]
        data_overview_slide = prs.slides.add_slide(data_overview_layout)
        title = data_overview_slide.shapes.title
        content = data_overview_slide.placeholders[1]
        title.text = "Data Overview"
        content.text = (
            "The data overview slide provides a summary of the key data points collected "
            "from various sources across the London and South regions."
        )

        # Slide 4: Example Chart
        example_chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = example_chart_slide.shapes.title
        title.text = "Example Chart"

        # Add a chart or image (assuming you have a chart image saved as 'chart.png')
        # For this example, let's assume we have an image stream from a chart

        # Example data for a chart
        data = {'Region': ['Central London', 'North West'], 'Value': [75, 65]}
        df = pd.DataFrame(data)

        # Generate a simple bar plot
        import matplotlib.pyplot as plt
        from io import BytesIO

        fig, ax = plt.subplots()
        df.plot(kind='bar', x='Region', y='Value', ax=ax, legend=False)
        ax.set_ylabel('Value')
        ax.set_title('Example Chart')

        # Save plot to a bytes buffer
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        img_stream.seek(0)

        # Add image to slide
        left = Inches(2)
        top = Inches(1.5)
        pic = example_chart_slide.shapes.add_picture(img_stream, left, top, width=Inches(6), height=Inches(4))

        # Save the presentation
        prs.save('London_and_South_Data_Report.pptx')

        print("PowerPoint report generated successfully!")

    except Exception as e:
        print(f"An error occurred: {str(e)}")


# Generate the PowerPoint report
create_powerpoint_report()

# Function to create PowerPoint report

# Initial setup of the application
create_background()
add_corner_image()
add_info_label()

# Labels and login button for initial login page
label = Label(window, text="Welcome to the Report Automator Application!", font=('Amasis MT Pro Black', 30, 'bold'),
              background='orange', bd=10)
label.place(x=40, y=200)

button = Button(window, text="Login to Start!", command=login, font=('Amasis MT Pro Black', 30, 'bold'), fg='green',
                background='orange', relief=SUNKEN, bd=10)
button.place(x=250, y=285)

button = Button(window, text="EXIT", command=exit_app, font=('Amasis MT Pro Black', 20, 'bold'), fg='red',
                background='orange', relief=SUNKEN, bd=10)
button.place(x=755, y=500)

# Start the main loop of the application
window.mainloop()
